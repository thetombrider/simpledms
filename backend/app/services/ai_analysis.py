from typing import Dict, List, Tuple
from anthropic import AsyncAnthropic, APIError, APIStatusError, RateLimitError, AuthenticationError
from ..core.config import settings
import PyPDF2
import io
import mimetypes
from docx.api import Document
from pptx import Presentation
import json
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class AIServiceError(Exception):
    """Base class for AI service errors"""
    def __init__(self, message: str, error_type: str):
        self.message = message
        self.error_type = error_type
        super().__init__(self.message)
        # Log error details
        logger.error(f"AIServiceError: {error_type} - {message}")

class AIAnalysisService:
    def __init__(self):
        self.api_key = settings.ANTHROPIC_API_KEY
        self.is_available = bool(self.api_key)
        if self.is_available:
            self.client = AsyncAnthropic(api_key=self.api_key)
            logger.info("AIAnalysisService initialized successfully")
        else:
            logger.warning("AIAnalysisService initialized without API key")
    
    def _check_api_key(self):
        """Check if Anthropic API key is configured"""
        if not self.is_available:
            logger.error("Anthropic API key not configured")
            raise AIServiceError(
                "Anthropic integration is not available - ANTHROPIC_API_KEY not configured",
                "configuration_error"
            )
    
    def _extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PowerPoint file"""
        prs = Presentation(io.BytesIO(file_content))
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        
        text = "\n".join(filter(None, text_parts))  # Filter out empty strings
        logger.info(f"Successfully extracted {len(text)} characters from PowerPoint")
        return text
    
    async def extract_text(self, file_content: bytes, mime_type: str) -> str:
        """Extract text content from various file types"""
        try:
            logger.info(f"Extracting text from file with mime type: {mime_type}")
            if mime_type == "application/pdf":
                # Handle PDF files
                pdf_file = io.BytesIO(file_content)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                logger.info(f"Successfully extracted {len(text)} characters from PDF")
                return text
            
            elif mime_type in ["application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
                # Handle Word documents
                doc = Document(io.BytesIO(file_content))
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                logger.info(f"Successfully extracted {len(text)} characters from Word document")
                return text
            
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                # Handle PowerPoint documents
                return self._extract_text_from_pptx(file_content)
            
            elif mime_type.startswith("text/"):
                # Handle text files
                text = file_content.decode('utf-8')
                logger.info(f"Successfully extracted {len(text)} characters from text file")
                return text
            
            else:
                supported_types = [
                    "application/pdf",
                    "application/msword",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "text/*"
                ]
                logger.error(f"Unsupported file type: {mime_type}. Supported types: {supported_types}")
                raise ValueError(f"Unsupported file type: {mime_type}. Supported types are: PDF, Word, PowerPoint, and text files.")
        
        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}", exc_info=True)
            raise ValueError(f"Error extracting text: {str(e)}")

    async def analyze_document(self, file_content: bytes, mime_type: str) -> Tuple[str, List[str], List[str]]:
        """
        Analyze document content and return:
        - Summary description
        - Suggested categories
        - Suggested tags
        """
        try:
            self._check_api_key()
            
            # Extract text from document
            text = await self.extract_text(file_content, mime_type)
            
            # Truncate text if too long (Claude has token limits)
            max_chars = 10000  # Claude can handle more tokens than GPT-3.5
            if len(text) > max_chars:
                logger.warning(f"Text truncated from {len(text)} to {max_chars} characters")
                text = text[:max_chars] + "..."
            
            # Create prompt for Claude
            system_prompt = """You are a document analysis assistant. Your task is to analyze documents and provide:
1. A brief summary (2-3 sentences)
2. Suggested categories (choose from: Invoice, Contract, Report, Other)
3. Relevant tags (3-5 keywords)

Always respond in valid JSON format with the following structure:
{
    "summary": "brief summary here",
    "categories": ["category1", "category2"],
    "tags": ["tag1", "tag2", "tag3"]
}"""

            user_prompt = f"""Please analyze this document:

{text}"""

            try:
                logger.info("Sending request to Anthropic API")
                # Get AI analysis using Claude
                response = await self.client.messages.create(
                    model="claude-3-sonnet-20240229",
                    max_tokens=1000,
                    temperature=0.3,
                    system=system_prompt,
                    messages=[
                        {"role": "user", "content": user_prompt}
                    ]
                )
                
                logger.info("Received response from Anthropic API")
                # Parse JSON from Claude's response
                content = response.content[0].text
                logger.debug(f"Raw API response: {content}")
                
                # Remove any markdown code block markers if present
                content = content.replace("```json", "").replace("```", "").strip()
                logger.debug(f"Cleaned response content: {content}")
                
                try:
                    result = json.loads(content)
                except json.JSONDecodeError as e:
                    logger.error(f"JSON parsing error. Content: {content}", exc_info=True)
                    raise AIServiceError(
                        f"Failed to parse AI response as JSON: {str(e)}. Response: {content}",
                        "processing_error"
                    )
                
                # Ensure we have all required fields
                if not all(k in result for k in ["summary", "categories", "tags"]):
                    missing_fields = [k for k in ["summary", "categories", "tags"] if k not in result]
                    logger.error(f"Missing required fields in response: {missing_fields}. Response: {result}")
                    raise AIServiceError(
                        f"Invalid response format from AI service. Missing fields: {missing_fields}",
                        "processing_error"
                    )
                
                logger.info("Successfully parsed AI response")
                return (
                    result["summary"],
                    result["categories"],
                    result["tags"]
                )
            
            except RateLimitError as e:
                logger.error("Rate limit exceeded", exc_info=True)
                raise AIServiceError(
                    "Anthropic service quota exceeded. Please check your billing details.",
                    "quota_exceeded"
                )
            except AuthenticationError:
                logger.error("Authentication failed", exc_info=True)
                raise AIServiceError(
                    "Invalid Anthropic API key. Please check your configuration.",
                    "authentication_error"
                )
            except APIStatusError as e:
                logger.error(f"API error: {str(e)}", exc_info=True)
                raise AIServiceError(
                    f"Anthropic service error: {str(e)}",
                    "api_error"
                )
            except Exception as e:
                logger.error(f"Unexpected error processing AI response: {str(e)}", exc_info=True)
                raise AIServiceError(
                    f"Error processing AI response: {str(e)}",
                    "processing_error"
                )
                
        except AIServiceError:
            raise
        except Exception as e:
            logger.error(f"Unexpected error in analyze_document: {str(e)}", exc_info=True)
            raise AIServiceError(f"Error analyzing document: {str(e)}", "unknown_error")

    async def get_summary(self, file_content: bytes, mime_type: str) -> str:
        """Get just the summary of a document"""
        summary, _, _ = await self.analyze_document(file_content, mime_type)
        return summary

    async def get_suggestions(self, file_content: bytes, mime_type: str) -> Dict[str, List[str]]:
        """Get category and tag suggestions"""
        _, categories, tags = await self.analyze_document(file_content, mime_type)
        return {
            "categories": categories,
            "tags": tags
        } 